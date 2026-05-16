"""
convertir.py — Conversión de formatos al vuelo (PDF, PPTX, DOCX, XLSX → Markdown).

Cambios respecto a la versión anterior:
  - limpiar_ruido_estatico y limpiar_markdown importadas desde utils.py.
    Ya no se definen localmente (eliminada la deuda técnica de duplicación).
"""

import pymupdf4llm
from markitdown import MarkItDown
from pptx import Presentation
import pandas as pd
import sys
import os
import time
import logging
import argparse

from utils import limpiar_ruido_estatico, limpiar_markdown  # fuente única de verdad

try:
    from tqdm import tqdm
    TQDM_DISPONIBLE = True
except ImportError:
    TQDM_DISPONIBLE = False

EXTENSIONES_VALIDAS = ['.pdf', '.pptx', '.docx', '.xlsx']
MAX_REINTENTOS = 2


# ---------------------------------------------------------------------------
# MOTORES ESPECÍFICOS POR FORMATO
# ---------------------------------------------------------------------------

def procesar_pptx_mejorado(ruta: str) -> str:
    """Estructura el PowerPoint por diapositivas e incluye notas del orador."""
    prs = Presentation(ruta)
    output = [f"# PRESENTACIÓN: {os.path.basename(ruta)}\n"]

    for i, slide in enumerate(prs.slides, 1):
        output.append(f"\n---\n## DIAPOSITIVA {i}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                output.append(shape.text.strip())

        if slide.has_notes_slide:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                output.append(f"\n> **Notas del Orador:** {notas}")

    return "\n".join(output)


def procesar_xlsx_mejorado(ruta: str) -> str:
    """Convierte cada hoja de Excel en una tabla Markdown independiente."""
    output = [f"# EXCEL: {os.path.basename(ruta)}\n"]

    with pd.ExcelFile(ruta) as xls:
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            output.append(f"\n## HOJA: {sheet_name}\n")
            if not df.empty:
                output.append(df.to_markdown(index=False))
            else:
                output.append("*Esta hoja está vacía*")

    return "\n".join(output)


# ---------------------------------------------------------------------------
# LÓGICA DE CONVERSIÓN
# ---------------------------------------------------------------------------

def ejecutar_conversion(
    ruta_archivo: str, logger: logging.Logger | None, sin_imagenes: bool = False
) -> str:
    """
    Convierte un archivo al formato Markdown y lo guarda en disco.

    Returns:
        "ok"      — conversión exitosa.
        "omitido" — el .md ya existía (caché).
        "error"   — falló tras MAX_REINTENTOS intentos.
    """
    ext = os.path.splitext(ruta_archivo)[1].lower()
    nombre_base = os.path.splitext(os.path.basename(ruta_archivo))[0]
    carpeta_destino = os.path.join(
        os.path.dirname(ruta_archivo), f"MD_{nombre_base}"
    )
    md_esperado = os.path.join(carpeta_destino, f"{nombre_base}.md")

    if os.path.isfile(md_esperado):
        return "omitido"

    os.makedirs(carpeta_destino, exist_ok=True)
    ultimo_error = None

    for _ in range(MAX_REINTENTOS):
        try:
            if ext == '.pdf':
                dir_actual = os.getcwd()
                os.chdir(carpeta_destino)
                try:
                    md_text = pymupdf4llm.to_markdown(
                        ruta_archivo,
                        write_images=not sin_imagenes,
                        margins=(50, 0, 50, 0),
                    )
                finally:
                    os.chdir(dir_actual)

            elif ext == '.pptx':
                md_text = procesar_pptx_mejorado(ruta_archivo)

            elif ext == '.xlsx':
                md_text = procesar_xlsx_mejorado(ruta_archivo)

            else:  # .docx
                md_engine = MarkItDown()
                md_text = md_engine.convert(ruta_archivo).text_content

            # Limpieza tipográfica ligera (centralizada en utils.py)
            contenido_final = limpiar_markdown(md_text)

            with open(md_esperado, 'w', encoding='utf-8') as f:
                f.write(contenido_final)

            return "ok"

        except Exception as e:
            ultimo_error = e
            time.sleep(0.5)

    if logger:
        logger.error(f"FALLO en {nombre_base}: {ultimo_error}")
    else:
        print(f"[!] Error en {nombre_base}: {ultimo_error}")

    return "error"


# ---------------------------------------------------------------------------
# INTERFAZ CLI INDEPENDIENTE
# ---------------------------------------------------------------------------

def resolver_ruta_inteligente(entrada: str) -> str | None:
    """Resuelve la ruta si el usuario olvida la extensión."""
    if os.path.isfile(entrada):
        return os.path.abspath(entrada)
    for ext in EXTENSIONES_VALIDAS:
        candidato = entrada + ext
        if os.path.isfile(candidato):
            return os.path.abspath(candidato)
    return None


def main() -> None:
    parser = argparse.ArgumentParser(prog="cvt", description="Conversor Universidad Pro")
    parser.add_argument("entrada", help="Archivo o carpeta")
    parser.add_argument("--quiet", "-q", action="store_true", help="Modo silencioso")
    args = parser.parse_args()

    ruta_abs = os.path.abspath(args.entrada)
    log_dir = ruta_abs if os.path.isdir(ruta_abs) else os.path.dirname(ruta_abs)

    logging.basicConfig(
        filename=os.path.join(log_dir, "errores.log"),
        level=logging.ERROR,
        format="%(asctime)s - %(message)s",
    )
    logger = logging.getLogger()
    contadores = {"ok": 0, "omitido": 0, "error": 0}

    if os.path.isdir(ruta_abs):
        archivos = [
            os.path.join(ruta_abs, f)
            for f in os.listdir(ruta_abs)
            if os.path.splitext(f)[1].lower() in EXTENSIONES_VALIDAS
        ]

        if not archivos:
            print("[!] No hay archivos compatibles.")
            return

        iterador = (
            tqdm(archivos, desc="Convirtiendo biblioteca", unit="fich", ncols=80)
            if TQDM_DISPONIBLE and not args.quiet
            else archivos
        )

        for r in iterador:
            res = ejecutar_conversion(r, logger, sin_imagenes=False)
            contadores[res] += 1

    else:
        ruta = resolver_ruta_inteligente(args.entrada)
        if not ruta:
            print(f"[!] Error: No se encontró '{args.entrada}'")
            return
        res = ejecutar_conversion(ruta, logger, sin_imagenes=False)
        contadores[res] += 1

    if not args.quiet:
        print("\n" + "═" * 35)
        print(f" PROCESO COMPLETADO")
        print(f" ✓ Convertidos: {contadores['ok']}")
        print(f" ↷ Ignorados:   {contadores['omitido']}")
        print(f" ✗ Fallidos:    {contadores['error']}")
        print("═" * 35)
        if contadores["error"] > 0:
            print(f"Detalles en: {os.path.join(log_dir, 'errores.log')}")


if __name__ == "__main__":
    main()
